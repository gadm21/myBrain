#!/usr/bin/env python3

"""
Reference Materials Processing Module

This module handles reading and processing reference materials for the AI Agent.
It supports various file formats including text files, PDFs, CSVs, PowerPoint presentations, and images.
"""

import io
import logging
import os
from typing import Any, Dict, List, Optional

import pytesseract
from PIL import Image

# Try importing optional dependencies
try:
    import pdfplumber

    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    import pandas as pd

    CSV_SUPPORT = True
except ImportError:
    CSV_SUPPORT = False

try:
    from pptx import Presentation

    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

import aiagent.memory as memory





def read_references(limit: int = 12000, file_paths: Optional[List[str]] = None) -> Dict[str, Any]:
    """Read and process reference files.

    Args:
        limit: Maximum number of characters to process per file
        file_paths: Optional list of file paths to process. If provided, only these files will be processed and returned.

    Returns:
        Dictionary containing processed reference data

    :noindex:
    """
    references = {}
    text_extensions = (".txt", ".md")
    max_chars = limit  # Limit to ~3000 tokens to leave room for other context
    total_files_processed = 0
    total_content_chars = 0

    REFERENCES_DIR = memory.REFERENCES_DIR
    
    logging.info(f"Reading references from directory: {REFERENCES_DIR}")
    files = [] 
    
    # read filepaths or REFERENCES_DIR
    if not file_paths:
        # Check if the references directory exists
        if not os.path.exists(REFERENCES_DIR):
            logging.warning(f"References directory does not exist: {REFERENCES_DIR}")
            # Check if we're in a serverless environment
            if os.environ.get("VERCEL"):
                logging.info("Running in serverless environment with read-only filesystem")
                # Just return empty references, don't try to create directories
                # /tmp directory paths should already be set up in memory/__init__.py
                return references
            else:
                # Try to create it in regular environments
                try:
                    os.makedirs(REFERENCES_DIR, exist_ok=True)
                    logging.info(f"Created references directory: {REFERENCES_DIR}")
                except Exception as e:
                    logging.error(f"Failed to create references directory: {e}")
                return references  # Return empty references
            
        # Check if the directory is accessible
        try:
            files = [os.path.join(REFERENCES_DIR, f) for f in os.listdir(REFERENCES_DIR)]
        except Exception as e:
            logging.error(f"Error accessing references directory: {e}")
            return references  # Return empty references

    else:
        files = file_paths
        
    
    for file_path in files:
        filename = os.path.basename(file_path)
        if not os.path.isfile(file_path):
            logging.debug(f"Skipping non-file item: {filename}")
            continue

        file_ext = os.path.splitext(filename)[1].lower()
        print("file_ext")
        print(file_ext)
        logging.debug(f"Processing file: {filename} (type: {file_ext})")
        try:
            content = ""
            processed = False  # Flag to check if file type was handled
            if file_ext in text_extensions:
                # Handle text files
                logging.debug(f"Reading text file: {filename}")
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                processed = True
            elif file_ext == ".py":
                # Handle Python files
                logging.debug(f"Reading Python file: {filename}")
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                processed = True
            elif file_ext == ".html":
                # Handle HTML files
                logging.debug(f"Reading HTML file: {filename}")
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                processed = True
            elif file_ext == ".pdf" and PDF_SUPPORT:
                # Handle PDF files
                logging.debug(f"Reading PDF file: {filename}")
                content = "PDF CONTENT:\n\n"
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        content += f"--- Page {i+1} ---\n{page_text}\n\n"
                processed = True
            elif file_ext == ".csv" and CSV_SUPPORT:
                # Handle CSV files
                logging.debug(f"Reading CSV file: {filename}")
                df = pd.read_csv(file_path)
                content = f"CSV DATA ({len(df)} rows, {len(df.columns)} columns):\n\n"
                # Add column headers
                content += "Columns: " + ", ".join(df.columns) + "\n\n"
                # Add sample data (first 10 rows)
                content += df.head(10).to_string(index=False) + "\n"
                if len(df) > 10:
                    content += f"\n... {len(df) - 10} more rows ...\n"
                processed = True
            elif file_ext == ".pptx" and PPTX_SUPPORT:
                # Handle PowerPoint files
                logging.debug(f"Reading PowerPoint file: {filename}")
                prs = Presentation(file_path)
                content = "POWERPOINT CONTENT:\n\n"
                for i, slide in enumerate(prs.slides):
                    content += f"--- Slide {i+1} ---\n"
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                content += paragraph.text + "\n"
                    content += "\n"
                processed = True
            elif file_ext in (".png", ".jpg", ".jpeg"):
                # Handle image files with OCR
                logging.debug(f"Processing image file with OCR: {filename}")
                try:
                    image = Image.open(file_path)
                    content = "IMAGE CONTENT (OCR):\n\n"
                    content += (
                        pytesseract.image_to_string(image)
                        or "No text detected in image."
                    )
                    processed = True
                except Exception as img_e:
                    logging.error(
                        f"Error processing image {filename} with OCR: {img_e}"
                    )

            # Add to references if processed successfully and has content
            if processed and content:
                # Truncate if exceeds max_chars
                if len(content) > max_chars:
                    logging.warning(
                        f"Truncating {filename} content from {len(content)} to {max_chars} chars"
                    )
                    content = (
                        content[:max_chars] + "\n... [content truncated due to length]"
                    )

                references[file_path] = content
                total_files_processed += 1
                total_content_chars += len(content)
            else:
                # If not processed, it means we don't support this file type
                if not processed:
                    logging.debug(f"Skipping unsupported file type: {filename}")
        except Exception as e:
            logging.error(f"Error processing file {filename}: {e}")

    logging.info(
        f"Processed {total_files_processed} reference files with {total_content_chars} total characters"
    )
    return references
